#!/usr/bin/env python3
"""
Kitchmatics Presentation PPTX Generator - Full Version (75 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
MAIN_COLOR = RGBColor(44, 62, 80)
ACCENT_COLOR = RGBColor(231, 76, 60)
SECONDARY_COLOR = RGBColor(52, 152, 219)
SUCCESS_COLOR = RGBColor(39, 174, 96)
WARNING_COLOR = RGBColor(243, 156, 18)
LIGHT_BG = RGBColor(236, 240, 241)
WHITE = RGBColor(255, 255, 255)

class PPTXGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle, tags=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = MAIN_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = SECONDARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        return slide

    def add_section_slide(self, part_num, title, subtitle):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background gradient effect (solid color)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = MAIN_COLOR
        bg.line.fill.background()

        # Part number
        part_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.8))
        tf = part_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Part {part_num}"
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(self, title, content_items, notes=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = MAIN_COLOR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if item.startswith('##'):
                p.text = item.replace('##', '').strip()
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = SECONDARY_COLOR
                p.space_before = Pt(20)
            elif item.startswith('-'):
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR
                p.space_before = Pt(8)
                p.level = 0
            elif item.startswith('  -'):
                p.text = "  ◦ " + item[3:].strip()
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(100, 100, 100)
                p.space_before = Pt(4)
                p.level = 1
            elif item == '':
                p.text = ''
                p.space_before = Pt(10)
            else:
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR

        # Notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def add_two_column_slide(self, title, left_title, left_items, right_title, right_items, notes=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = MAIN_COLOR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Left column title
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR

        # Left content
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = MAIN_COLOR
            p.space_before = Pt(8)

        # Right column title
        right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR

        # Right content
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = MAIN_COLOR
            p.space_before = Pt(8)

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def add_problem_slide(self, title, items, notes=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header with warning color
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = ACCENT_COLOR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "⚠️ " + title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if item.startswith('##'):
                p.text = item.replace('##', '').strip()
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = ACCENT_COLOR
                p.space_before = Pt(20)
            elif item.startswith('-'):
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR
                p.space_before = Pt(8)
            else:
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_solution_slide(self, title, items, notes=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header with success color
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = SUCCESS_COLOR
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if item.startswith('##'):
                p.text = item.replace('##', '').strip()
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = SUCCESS_COLOR
                p.space_before = Pt(20)
            elif item.startswith('-'):
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR
                p.space_before = Pt(8)
            else:
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = MAIN_COLOR

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def generate(self):
        # ==================== PART 1: INTRODUCTION ====================

        # Slide 1: Title
        self.add_title_slide(
            "🍔 Kitchmatics",
            "AI-Powered Autonomous Restaurant System"
        )

        # Slide 2: Agenda
        self.add_content_slide(
            "📋 Agenda",
            [
                "## Part 1-3",
                "- 프로젝트 소개",
                "- 시스템 아키텍처",
                "- FMS 핵심 설계",
                "",
                "## Part 4-5",
                "- 모바일 로봇 시스템",
                "- 문제 해결 사례 (4건)",
                "",
                "## Part 6-11",
                "- AI 서버 통합",
                "- 로봇팔, GUI, 테스트, 성능, 결론"
            ],
            "오늘 발표는 총 11개 파트로 구성되어 있습니다."
        )

        # Slide 3: Team Introduction
        self.add_content_slide(
            "👥 팀 소개",
            [
                "## 팀 구성",
                "- 6명 팀원, 8주 개발",
                "",
                "## 역할 분담",
                "- FMS/Backend: 중앙 관제 시스템, 데이터베이스",
                "- Navigation: 로봇 네비게이션, SLAM, 경로계획",
                "- AI Integration: YOLO, 음성인식, LLM 연동",
                "- HW/Embedded: 로봇팔 제어, 센서 통합"
            ],
            "6명의 팀원이 8주간 개발했습니다."
        )

        # Slide 4: Project Background
        self.add_content_slide(
            "🎯 프로젝트 배경",
            [
                "## 문제 인식",
                "- 외식업계 인력난 심화",
                "- 서빙 로봇 도입 증가 추세",
                "- 단순 서빙을 넘어선 통합 솔루션 필요",
                "",
                "## 우리의 접근",
                "- 주문 → 조리 → 서빙 전 과정 자동화",
                "- 다중 로봇 협업으로 효율성 극대화",
                "- AI 기반 품질 관리 및 음성 주문"
            ],
            "외식업계의 인력난 문제를 해결하기 위해 전체 운영 자동화를 목표로 했습니다."
        )

        # Slide 5: Project Goals
        self.add_content_slide(
            "🎯 프로젝트 목표",
            [
                "## 1. 다중 로봇 관제",
                "- 3대 이상의 서빙 로봇을 중앙에서 효율적으로 관리",
                "",
                "## 2. End-to-End 자동화",
                "- 주문 접수부터 음식 전달까지 무인 운영",
                "",
                "## 3. AI 품질 관리",
                "- 컴퓨터 비전으로 조리 완성도 검증",
                "",
                "## 4. 자연어 인터페이스",
                "- 음성 명령으로 직관적인 주문 경험"
            ],
            "4가지 핵심 목표를 설정했습니다."
        )

        # Slide 6: Key Metrics
        self.add_content_slide(
            "📊 프로젝트 지표",
            [
                "## 코드베이스",
                "- 10,225+ Lines of Code",
                "- 155 Test Cases",
                "- 23 ROS2 Nodes",
                "",
                "## 하드웨어",
                "- 3 Mobile Robots (TurtleBot3)",
                "- 2 Robot Arms (MyCobot 280)",
                "- 5 Domain IDs"
            ],
            "프로젝트 규모를 숫자로 보여드립니다."
        )

        # Slide 7: Technology Stack
        self.add_two_column_slide(
            "🛠 기술 스택",
            "Core Platform",
            [
                "ROS2 Humble - 로봇 미들웨어",
                "CycloneDDS - DDS 구현체",
                "Nav2 - 자율 주행",
                "Python 3.10+ - 주 개발 언어",
                "YOLOv8 - 객체 탐지",
                "OpenAI APIs - 음성/자연어"
            ],
            "Backend & Frontend",
            [
                "FastAPI - REST API",
                "PostgreSQL - 데이터베이스",
                "SQLAlchemy - ORM",
                "PyQt5 - Desktop GUI",
                "TCP/JSON - 실시간 통신",
                "pytest - 테스트 프레임워크"
            ],
            "ROS2 Humble을 핵심으로 다양한 기술을 통합했습니다."
        )

        # ==================== PART 2: SYSTEM ARCHITECTURE ====================

        # Slide 8: Section Divider
        self.add_section_slide(2, "System Architecture", "시스템 아키텍처")

        # Slide 9: High-Level Architecture
        self.add_content_slide(
            "🏗 전체 시스템 구성",
            [
                "## 핵심 구성요소",
                "- Customer GUI (Kiosk) - 고객 주문 인터페이스",
                "- Admin GUI (Control) - 관리자 모니터링",
                "- AI Server (YOLO/Voice) - 컴퓨터 비전 및 음성 처리",
                "- FMS (Fleet Manager) - 중앙 관제 시스템",
                "- Mobile Robots (pinky1, pinky2, pinky3) - 서빙 로봇",
                "- Robot Arms - 조리 로봇팔",
                "",
                "## 통신 방식",
                "- GUI ↔ FMS: TCP/JSON",
                "- FMS ↔ Robots: ROS2 DDS"
            ],
            "FMS가 중앙에서 모든 컴포넌트를 관제합니다."
        )

        # Slide 10: Hardware Components
        self.add_two_column_slide(
            "🤖 하드웨어 구성",
            "Mobile Robots (3대)",
            [
                "Platform: TurtleBot3 Burger",
                "Sensors: LiDAR, IMU, Encoders",
                "Computer: Raspberry Pi 4",
                "Payload: 음식 트레이"
            ],
            "Robot Arms (2대)",
            [
                "Model: MyCobot 280",
                "DOF: 6축",
                "Task: 샌드위치 조립",
                "Interface: Python SDK"
            ],
            "TurtleBot3 3대와 MyCobot 2대를 사용했습니다."
        )

        # Slide 11: Network Topology
        self.add_content_slide(
            "🌐 네트워크 토폴로지",
            [
                "## Master PC (192.168.1.3)",
                "- FMS Node, AI Server 실행",
                "- GUI와 TCP:9000으로 연결",
                "",
                "## Mobile Robots",
                "- pinky1: 192.168.1.7 (Domain 11)",
                "- pinky2: 192.168.1.6 (Domain 12)",
                "- pinky3: 192.168.1.11 (Domain 13)",
                "",
                "## Robot Arm",
                "- 192.168.1.4 (Domain 20)"
            ],
            "GUI는 TCP, 로봇들은 ROS2 DDS로 통신합니다."
        )

        # Slide 12: ROS2 Domain Architecture
        self.add_content_slide(
            "🔗 ROS2 Multi-Domain 설계",
            [
                "## Domain ID 배정",
                "- Domain 25: FMS (Master) - 중앙 관제, Domain Bridge",
                "- Domain 11: pinky1 - 서빙 로봇 #1",
                "- Domain 12: pinky2 - 서빙 로봇 #2",
                "- Domain 13: pinky3 - 서빙 로봇 #3",
                "- Domain 20: Robot Arms - 조리 로봇팔",
                "",
                "## Why Multi-Domain?",
                "- 각 로봇이 동일한 토픽명(/cmd_vel, /odom) 사용",
                "- Domain으로 분리하여 충돌 방지"
            ],
            "5개의 Domain ID를 사용합니다."
        )

        # Slide 13: Domain Bridge
        self.add_content_slide(
            "🌉 Domain Bridge 구조",
            [
                "## 토픽 리매핑",
                "- /cmd_vel (Domain 11) → /pinky1/cmd_vel (Domain 25)",
                "- /odom (Domain 11) → /pinky1/odom (Domain 25)",
                "- /amcl_pose (Domain 11) → /pinky1/amcl_pose (Domain 25)",
                "",
                "## 설정 예시",
                "- domain_bridge.yaml에서 매핑 정의",
                "- FMS에서 /pinky1/cmd_vel로 명령",
                "- pinky1에서는 /cmd_vel로 수신"
            ],
            "Domain Bridge가 토픽을 리매핑합니다."
        )

        # Slide 14: Clean Architecture
        self.add_content_slide(
            "🏛 Clean Architecture 적용",
            [
                "## 레이어 구조",
                "- Presentation Layer: FMS Node, GUI Controllers",
                "- Application Layer: OrderHandler, FleetController, UseCases",
                "- Domain Layer: Order, Robot, OrderWorkflow (Entities)",
                "- Infrastructure Layer: ROS2 Publishers, TCP Server, Database",
                "",
                "## 의존성 규칙",
                "- 안쪽 레이어는 바깥쪽 레이어를 알지 못함",
                "- Domain Layer는 ROS2에 의존하지 않음"
            ],
            "Clean Architecture로 레이어별 책임을 분리했습니다."
        )

        # ==================== PART 3: FMS DEEP DIVE ====================

        # Slide 15: Section Divider
        self.add_section_slide(3, "FMS Deep Dive", "Fleet Management System 심층 분석")

        # Slide 16: FMS Overview
        self.add_two_column_slide(
            "🎛 FMS 개요",
            "핵심 기능",
            [
                "주문 수신 및 처리",
                "로봇 배차 및 경로 지정",
                "조리-서빙 동기화",
                "실시간 상태 모니터링",
                "에러 감지 및 복구"
            ],
            "주요 모듈",
            [
                "fms_node.py - 메인 노드",
                "order_handler.py - 주문 처리",
                "fleet_controller.py - 로봇 관리",
                "error_detector.py - 에러 감지",
                "gui_tcp_server.py - GUI 통신"
            ],
            "FMS는 전체 시스템의 두뇌입니다."
        )

        # Slide 17: FMS Node Structure
        self.add_content_slide(
            "📦 FMS Node 구조",
            [
                "## 주요 컴포넌트",
                "- Application Layer: OrderHandler, FleetController, ErrorDetector",
                "- Infrastructure Layer: GUITCPServer (port=9000)",
                "- ROS2 Communication: Navigation Action Clients, Publishers",
                "",
                "## Dependency Injection 패턴",
                "- OrderHandler에 콜백 등록",
                "- send_cooking_command, navigate_robot, send_gui_notification",
                "- Application Layer가 Infrastructure에 직접 의존하지 않음"
            ],
            "FMS Node는 Dependency Injection 패턴을 사용합니다."
        )

        # Slide 18: Order Handler
        self.add_content_slide(
            "📋 Order Handler",
            [
                "## 책임",
                "- 주문 Workflow 관리",
                "- 상태 전이 제어",
                "- 콜백 트리거",
                "",
                "## 상태 전이",
                "- RECEIVED → COOKING → LOADING",
                "- LOADING → LOADED → DELIVERING",
                "- DELIVERING → ARRIVED → COMPLETED"
            ],
            "Order Handler는 주문의 전체 생명주기를 관리합니다."
        )

        # Slide 19: Fleet Controller
        self.add_content_slide(
            "🚗 Fleet Controller",
            [
                "## 로봇 상태",
                "- IDLE: 대기 중",
                "- MOVING: 이동 중",
                "- BUSY: 배달 중",
                "- ERROR: 에러 발생",
                "",
                "## 배차 알고리즘",
                "- 우선순위: IDLE 상태 로봇 선택",
                "- 거리 기반: 주방과 가장 가까운 로봇 배정"
            ],
            "Fleet Controller가 로봇 배차를 담당합니다."
        )

        # Slide 20: Order State Machine
        self.add_content_slide(
            "🔄 Order State Machine",
            [
                "## 상태 전이 흐름",
                "- RECEIVED: 주문 접수",
                "- COOKING: 조리 중 + 로봇 이동 중",
                "- LOADING: 조리 완료 대기",
                "- LOADED: 적재 완료",
                "- DELIVERING: 테이블로 이동 중",
                "- ARRIVED: 도착, 고객 확인 대기",
                "- COMPLETED: 완료"
            ],
            "각 상태 전이마다 특정 조건과 액션이 정의되어 있습니다."
        )

        # Slide 21: Error Detector
        self.add_two_column_slide(
            "⚠️ Error Detector",
            "감지 항목",
            [
                "Navigation Timeout: 목적지 미도착",
                "Robot Stuck: 이동 중 정지",
                "Communication Loss: 통신 두절",
                "Cooking Timeout: 조리 지연"
            ],
            "복구 옵션",
            [
                "RETRY: 동일 작업 재시도",
                "SKIP: 다음 단계로 진행",
                "MANUAL: 운영자 개입 요청"
            ],
            "Error Detector가 다양한 에러 상황을 감지합니다."
        )

        # ==================== PART 4: MOBILE ROBOTS ====================

        # Slide 22: Section Divider
        self.add_section_slide(4, "Mobile Robots", "모바일 로봇 시스템")

        # Slide 23: Navigation Stack
        self.add_content_slide(
            "🧭 Navigation Stack",
            [
                "## Nav2 구성요소",
                "- Nav2 BT Navigator: Behavior Tree 기반 의사결정",
                "- Planner (NavFn/Smac): 전역 경로 계획",
                "- Controller (DWB): 지역 경로 추종",
                "- Recovery Behaviors: 에러 복구 동작",
                "- Costmap 2D: 장애물 맵",
                "",
                "## 사용 설정",
                "- Global Planner: NavFn",
                "- Local Controller: DWB"
            ],
            "Nav2의 구조입니다."
        )

        # Slide 24: AMCL Localization
        self.add_two_column_slide(
            "📍 AMCL 위치 추정",
            "특징",
            [
                "Particle Filter 기반",
                "LiDAR + Odometry 융합",
                "맵 매칭으로 위치 보정"
            ],
            "주요 파라미터",
            [
                "min_particles: 500",
                "max_particles: 2000",
                "update_min_d: 0.2m",
                "update_min_a: 0.5rad"
            ],
            "AMCL이 파티클 필터로 로봇 위치를 추정합니다."
        )

        # Slide 25: Waypoint System
        self.add_content_slide(
            "📌 Waypoint 관리",
            [
                "## 정의된 Waypoints",
                "- kitchen (point13): 음식 픽업 위치",
                "- table1 ~ table4: 테이블 서빙 위치",
                "- pinky1_spot, pinky2_spot, pinky3_spot: 대기 위치",
                "",
                "## 설정 방식",
                "- waypoints.yaml 파일에서 좌표 관리",
                "- (x, y, θ) 형식으로 위치 및 방향 지정"
            ],
            "미리 정의된 Waypoint로 로봇을 이동시킵니다."
        )

        # ==================== PART 5: PROBLEMS & SOLUTIONS ====================

        # Slide 26: Section Divider
        self.add_section_slide(5, "Problems & Solutions", "문제 해결 사례")

        # Slide 27: Problem 1 - Symptom
        self.add_problem_slide(
            "Problem #1: DDS Discovery 실패",
            [
                "## 증상",
                "- Master PC에서 ros2 topic list 실행",
                "- 로봇의 토픽이 전혀 보이지 않음",
                "- SSH 연결은 정상 동작",
                "- ping도 성공",
                "",
                "## 이상한 점",
                "- 네트워크는 연결되어 있음",
                "- ROS2만 통신 안됨"
            ],
            "첫 번째 문제입니다. 로봇과 네트워크는 연결되어 있는데 ROS2 토픽이 보이지 않았습니다."
        )

        # Slide 28: Problem 1 - Investigation
        self.add_content_slide(
            "Problem #1: 조사 과정",
            [
                "## 디버깅 단계",
                "- 1. 로봇에서 직접 확인 → 로컬에서는 토픽 정상",
                "- 2. Domain ID 확인 → 양쪽 모두 같은 ID 사용",
                "- 3. Wireshark 패킷 분석 → Multicast 패킷 미수신!",
                "- 4. Router 설정 확인 → WiFi AP에서 Multicast 차단 발견",
                "",
                "## 발견",
                "- WiFi AP가 Multicast UDP를 차단하고 있었음"
            ],
            "단계별로 조사했습니다. 네트워크 레벨에서 문제를 발견했습니다."
        )

        # Slide 29: Problem 1 - Root Cause
        self.add_problem_slide(
            "Problem #1: 근본 원인",
            [
                "## ROS2 DDS 기본 동작",
                "- 노드 발견에 Multicast UDP 사용",
                "- 주소: 239.255.0.1",
                "- 포트: 7400-7500",
                "",
                "## 문제점",
                "- 대부분의 WiFi AP는 보안/성능상 Multicast 차단 또는 제한",
                "- 가정/기업 WiFi 공유기에서 흔한 설정"
            ],
            "근본 원인은 WiFi AP의 Multicast 차단이었습니다."
        )

        # Slide 30: Solution 1 - Approach
        self.add_solution_slide(
            "Solution #1: 접근 방식",
            [
                "## Option 1: Router 설정 변경",
                "- Multicast 허용",
                "- ❌ 모든 환경에서 불가능",
                "- ❌ 보안 우려",
                "",
                "## Option 2: Unicast Discovery ✓",
                "- 직접 IP 주소로 통신",
                "- ✅ 모든 WiFi에서 동작",
                "- ✅ CycloneDDS 지원"
            ],
            "Unicast Discovery가 더 범용적인 해결책이었습니다."
        )

        # Slide 31: Solution 1 - Implementation
        self.add_solution_slide(
            "Solution #1: CycloneDDS 설정",
            [
                "## XML 설정 파일",
                "- cyclonedds_main.xml에 Peer IP 명시",
                "- <Peer address='192.168.1.7'/> (pinky1)",
                "- <Peer address='192.168.1.6'/> (pinky2)",
                "",
                "## 환경 변수",
                "- RMW_IMPLEMENTATION=rmw_cyclonedds_cpp",
                "- CYCLONEDDS_URI=file:///path/to/config.xml"
            ],
            "CycloneDDS XML 파일에 Peer IP 주소를 명시합니다."
        )

        # Slide 32: Solution 1 - Result
        self.add_solution_slide(
            "Solution #1: 결과",
            [
                "## Before (Multicast)",
                "- WiFi AP에서 차단",
                "- 토픽 발견 불가",
                "",
                "## After (Unicast)",
                "- 직접 IP-to-IP 통신",
                "- 모든 토픽 정상 발견",
                "",
                "## 성과",
                "- 모든 WiFi 환경에서 안정적인 ROS2 통신 확보"
            ],
            "Unicast로 변경 후 모든 토픽이 정상적으로 발견됩니다."
        )

        # Slide 33: Problem 2 - Symptom
        self.add_problem_slide(
            "Problem #2: Topic Name Collision",
            [
                "## 증상",
                "- 모든 로봇이 같은 토픽명 사용: /cmd_vel, /odom",
                "- 한 도메인에서 3대 로봇 운영 시 데이터 혼선",
                "- 어떤 /amcl_pose가 어떤 로봇인지 구분 불가",
                "",
                "## ros2 topic list 결과",
                "- /cmd_vel ← pinky1? pinky2? pinky3?",
                "- /odom ← 어떤 로봇의 odometry?"
            ],
            "TurtleBot3의 기본 토픽명이 모두 동일해서 구분이 불가능했습니다."
        )

        # Slide 34: Problem 2 - Impact
        self.add_problem_slide(
            "Problem #2: 영향",
            [
                "## 운영상 문제",
                "- ❌ 잘못된 로봇에 명령 전달",
                "- ❌ 위치 추정 오류",
                "- ❌ 충돌 위험",
                "- ❌ Fleet 관리 불가능",
                "",
                "## 기술적 한계",
                "- ❌ TurtleBot3 펌웨어 수정 어려움",
                "- ❌ Nav2 내부 토픽명 하드코딩"
            ],
            "토픽 충돌로 인해 Fleet 관리가 불가능했습니다."
        )

        # Slide 35: Solution 2 - Design
        self.add_solution_slide(
            "Solution #2: Multi-Domain Architecture",
            [
                "## 핵심 아이디어",
                "- 각 로봇을 별도의 ROS2 Domain에서 실행",
                "- Domain Bridge로 연결",
                "",
                "## 구조",
                "- Domain 11 (pinky1): /cmd_vel, /odom",
                "- Domain 12 (pinky2): /cmd_vel, /odom",
                "- Domain 25 (FMS): /pinky1/cmd_vel, /pinky2/cmd_vel"
            ],
            "각 로봇이 독립된 Domain에서 동작합니다."
        )

        # Slide 36: Solution 2 - Configuration
        self.add_solution_slide(
            "Solution #2: Domain Bridge 설정",
            [
                "## YAML 설정",
                "- from_domain: 11, to_domain: 25",
                "- topic: /cmd_vel → remap_to: /pinky1/cmd_vel",
                "- topic: /odom → remap_from: /pinky1/odom",
                "",
                "## 실행",
                "- ros2 run domain_bridge domain_bridge config.yaml"
            ],
            "YAML 파일에서 토픽 리매핑을 정의합니다."
        )

        # Slide 37: Solution 2 - Result
        self.add_solution_slide(
            "Solution #2: 결과",
            [
                "## FMS Domain에서 토픽 목록",
                "- /pinky1/cmd_vel, /pinky1/odom, /pinky1/amcl_pose",
                "- /pinky2/cmd_vel, /pinky2/odom, /pinky2/amcl_pose",
                "",
                "## 성과",
                "- ✅ 로봇별 명확한 토픽 구분",
                "- ✅ 독립적인 로봇 제어 가능",
                "- ✅ 로봇 펌웨어 수정 불필요",
                "- ✅ 확장 용이 (로봇 추가 시 Domain만 추가)"
            ],
            "각 로봇의 토픽이 명확히 구분됩니다."
        )

        # Slide 38: Problem 3 - Symptom
        self.add_problem_slide(
            "Problem #3: 조리-로봇 동기화",
            [
                "## Case A: 로봇 먼저 도착",
                "- 로봇이 주방에 도착",
                "- 음식 아직 조리 중",
                "- 로봇 대기 → 효율 저하",
                "",
                "## Case B: 음식 먼저 완성",
                "- 조리 완료",
                "- 로봇 아직 이동 중",
                "- 음식 방치 → 품질 저하"
            ],
            "로봇 도착과 조리 완료 시점이 일치하지 않아 비효율이 발생했습니다."
        )

        # Slide 39: Problem 3 - Analysis
        self.add_problem_slide(
            "Problem #3: 실패한 접근",
            [
                "## Approach 1: 타이머 기반",
                "- 조리 시작 후 30초 대기 후 로봇 호출",
                "- ❌ 조리 시간 가변적",
                "- ❌ 로봇 이동 시간 가변적",
                "",
                "## Approach 2: 폴링",
                "- 주기적으로 상태 확인",
                "- ❌ CPU 낭비",
                "- ❌ 반응 지연"
            ],
            "타이머와 폴링 방식 모두 효과적이지 않았습니다."
        )

        # Slide 40: Solution 3 - Design
        self.add_solution_slide(
            "Solution #3: Event-Driven State Machine",
            [
                "## 핵심 원칙",
                "- 조리 완료와 로봇 도착을 독립적 이벤트로 처리",
                "- AND 조건으로 다음 단계 진행",
                "",
                "## 구현",
                "- on_cooking_complete(): cooking_complete = True",
                "- on_robot_arrived_kitchen(): robot_at_kitchen = True",
                "- _check_loading_ready(): 둘 다 True면 LOADED로 전이"
            ],
            "Event-Driven 방식으로 변경했습니다."
        )

        # Slide 41: Solution 3 - State Diagram
        self.add_content_slide(
            "Solution #3: 상태 다이어그램",
            [
                "## COOKING 상태에서 분기",
                "- Event A: Robot Arrived at Kitchen",
                "- Event B: Cooking Done",
                "",
                "## AND 조건 체크",
                "- A AND B가 모두 만족되면",
                "- → LOADED 상태로 전이",
                "",
                "## 장점",
                "- 이벤트 순서에 관계없이 동작",
                "- Race Condition 완전 방지"
            ],
            "두 이벤트가 어떤 순서로 발생하든, 둘 다 만족되면 LOADED로 전이합니다."
        )

        # Slide 42: Solution 3 - Result
        self.add_solution_slide(
            "Solution #3: 결과",
            [
                "## 측정 지표",
                "- Race Conditions: 0",
                "- 동기화 성공률: 100%",
                "- 대기 시간 감소: -40%",
                "",
                "## 추가 이점",
                "- ✅ 명확한 상태 추적 가능",
                "- ✅ 디버깅 용이 (상태 로그)",
                "- ✅ 에러 복구 지점 명확",
                "- ✅ 테스트 용이 (상태별 단위 테스트)"
            ],
            "State Machine 도입으로 동기화 문제가 완전히 해결되었습니다."
        )

        # Slide 43: Problem 4 - Requirements
        self.add_problem_slide(
            "Problem #4: GUI-FMS 통신",
            [
                "## 기능적 요구사항",
                "- 주문 전송 (GUI → FMS)",
                "- 도착 알림 (FMS → GUI, Push)",
                "- 수령 확인 (GUI → FMS)",
                "",
                "## 제약사항",
                "- ❌ GUI 디바이스에 ROS2 설치 불가",
                "- ❌ 태블릿/키오스크 환경",
                "- ✅ TCP/IP 사용 가능"
            ],
            "GUI 디바이스에는 ROS2를 설치할 수 없어서 별도의 통신 방법이 필요했습니다."
        )

        # Slide 44: Solution 4 - Protocol Design
        self.add_solution_slide(
            "Solution #4: TCP/JSON Protocol",
            [
                "## 메시지 포맷",
                "- 4-byte length header (Big-endian)",
                "- JSON payload (UTF-8 encoded)",
                "",
                "## 메시지 타입",
                "- new_order: 고객 주문",
                "- delivery_notification: Push 알림",
                "- delivery_complete: 수령 확인"
            ],
            "4바이트 길이 헤더 + JSON 페이로드 구조로 설계했습니다."
        )

        # Slide 45: Solution 4 - Implementation
        self.add_solution_slide(
            "Solution #4: 서버 구현",
            [
                "## GUITCPServer 클래스",
                "- Thread-per-client 모델",
                "- Handler 패턴으로 메시지 처리",
                "- broadcast()로 Push 알림",
                "",
                "## 확장성",
                "- register_handler()로 새 메시지 타입 추가",
                "- Open/Closed Principle 적용"
            ],
            "Thread-per-client 모델로 구현했습니다."
        )

        # Slide 46: Solution 4 - Result
        self.add_solution_slide(
            "Solution #4: 결과",
            [
                "## 성능",
                "- 주문 처리 지연: <100ms",
                "- Push 알림 지연: <50ms",
                "- 메시지 전달률: 100%",
                "",
                "## 장점",
                "- ✅ ROS2 없이 실시간 통신",
                "- ✅ 다양한 클라이언트 지원",
                "- ✅ Handler 패턴으로 확장 용이",
                "- ✅ JSON으로 디버깅 용이"
            ],
            "TCP/JSON 프로토콜로 요구사항을 모두 충족했습니다."
        )

        # ==================== PART 6: AI INTEGRATION ====================

        # Slide 47: Section Divider
        self.add_section_slide(6, "AI Integration", "AI 서버 통합")

        # Slide 48: AI Server Overview
        self.add_two_column_slide(
            "🤖 AI 서버 구성",
            "YOLO Server",
            [
                "Model: YOLOv8",
                "Task: 음식 품질 검사",
                "API: FastAPI REST",
                "Port: 8001"
            ],
            "Voice Server",
            [
                "STT: OpenAI Whisper",
                "NLU: GPT-4 Function Calling",
                "TTS: OpenAI TTS",
                "Port: 8002"
            ],
            "두 개의 AI 서버가 있습니다."
        )

        # Slide 49: YOLO Training
        self.add_two_column_slide(
            "📸 YOLOv8: 학습 데이터",
            "데이터셋",
            [
                "빵 (상단): 500장",
                "빵 (하단): 500장",
                "패티: 450장",
                "치즈: 400장",
                "양상추: 380장",
                "토마토: 350장"
            ],
            "학습 설정",
            [
                "Base Model: YOLOv8n",
                "Epochs: 100",
                "Image Size: 640x640",
                "Augmentation: Mosaic, Mixup"
            ],
            "샌드위치 재료별로 학습 데이터를 구축했습니다."
        )

        # Slide 50: YOLO Integration
        self.add_content_slide(
            "📸 YOLOv8: 통합",
            [
                "## API Endpoint",
                "- POST /detect",
                "- 입력: 이미지 파일",
                "- 출력: detections 배열 (class, confidence, bbox)",
                "",
                "## FMS 연동",
                "- 카메라 이미지 → YOLO 서버 전송",
                "- 필요한 재료 모두 감지되면 조리 완료 판단"
            ],
            "REST API로 YOLO 서버와 통신합니다."
        )

        # Slide 51: Voice Pipeline
        self.add_content_slide(
            "🎤 음성 인식 파이프라인",
            [
                "## 처리 흐름",
                "- 1. Microphone → Audio 입력",
                "- 2. Whisper (STT) → Text 변환",
                "- 3. GPT-4 Function Calling → 구조화 데이터",
                "- 4. TTS → 음성 응답",
                "",
                "## 예시",
                "- 입력: '치즈버거 두 개 주세요'",
                "- 출력: {menu: '치즈버거', quantity: 2}"
            ],
            "음성 → 텍스트 → 구조화 데이터 → 음성 응답의 파이프라인입니다."
        )

        # Slide 52: GPT Function Calling
        self.add_content_slide(
            "🎤 GPT-4 Function Calling",
            [
                "## Function 정의",
                "- name: create_order",
                "- parameters: menu_name, quantity, options",
                "",
                "## 동작 방식",
                "- 사용자 발화를 GPT-4에 전달",
                "- GPT-4가 적절한 function과 인자 결정",
                "- 구조화된 JSON 반환"
            ],
            "GPT-4 Function Calling으로 자연어를 구조화된 데이터로 변환합니다."
        )

        # ==================== PART 7: ROBOT ARM ====================

        # Slide 53: Section Divider
        self.add_section_slide(7, "Robot Arm", "로봇팔 시스템")

        # Slide 54: Robot Arm Hardware
        self.add_two_column_slide(
            "🦾 로봇팔 하드웨어",
            "MyCobot 280",
            [
                "DOF: 6축",
                "Payload: 250g",
                "Reach: 280mm",
                "Repeatability: ±0.5mm",
                "Interface: Serial/Python SDK"
            ],
            "End Effector",
            [
                "Type: Gripper",
                "Range: 0-70mm",
                "Force: 조절 가능"
            ],
            "MyCobot 280 6축 로봇팔을 사용합니다."
        )

        # Slide 55: Robot Arm FMS Integration
        self.add_content_slide(
            "🦾 FMS 연동",
            [
                "## 토픽 인터페이스",
                "- /arm/command (FMS → Arm): 조리 명령",
                "- /arm/status (Arm → FMS): 상태 보고",
                "",
                "## 명령 형식 (JSON)",
                "- job_id: 작업 식별자",
                "- operation: START | PAUSE | RESUME | CANCEL",
                "- order: 주문 정보"
            ],
            "FMS와 로봇팔은 ROS2 토픽으로 통신합니다."
        )

        # Slide 56: Recipe Execution
        self.add_content_slide(
            "🦾 레시피 실행",
            [
                "## RecipeExecutor",
                "- YAML로 정의된 레시피 로드",
                "- step별로 순차 실행",
                "- 진행률 실시간 보고",
                "",
                "## 레시피 단계",
                "- 하단빵 배치 → 패티 → 치즈 → 양상추 → 상단빵"
            ],
            "YAML로 정의된 레시피를 단계별로 실행합니다."
        )

        # ==================== PART 8: GUI ====================

        # Slide 57: Section Divider
        self.add_section_slide(8, "GUI Development", "사용자 인터페이스")

        # Slide 58: Customer GUI
        self.add_two_column_slide(
            "📱 Customer GUI",
            "주요 화면",
            [
                "메뉴 선택: 터치 친화적 그리드",
                "장바구니: 수량 조절, 옵션",
                "결제: QR 코드 / 카드",
                "대기: 주문 상태 표시",
                "도착 알림: 팝업 + 사운드"
            ],
            "기술 스택",
            [
                "Framework: PyQt5",
                "Communication: TCP/JSON",
                "Assets: 메뉴 이미지",
                "Animation: Qt Animation"
            ],
            "Customer GUI는 키오스크용으로 설계되었습니다."
        )

        # Slide 59: Order Flow
        self.add_content_slide(
            "📱 주문 플로우",
            [
                "## 단계별 화면",
                "- 1. 메뉴 선택 → 햄버거, 치즈버거 등 터치",
                "- 2. 장바구니 → 수량 조절, 옵션 선택",
                "- 3. 결제 → 총액 확인, 결제 진행",
                "- 4. 주문 대기 → '조리중...' 표시",
                "- 5. 도착 알림 → Push 알림 팝업",
                "- 6. 수령 확인 → 버튼 클릭"
            ],
            "메뉴 선택부터 도착 알림까지의 전체 플로우입니다."
        )

        # Slide 60: Admin GUI
        self.add_two_column_slide(
            "🖥 Admin GUI",
            "모니터링 기능",
            [
                "Fleet 상태: 로봇 위치/상태",
                "주문 현황: 실시간 리스트",
                "에러 알림: 즉시 표시",
                "통계: 일별/시간별 분석"
            ],
            "제어 기능",
            [
                "수동 배차: 로봇 지정",
                "긴급 정지: 전체/개별",
                "복귀 명령: 홈 위치",
                "메뉴 관리: 품절 설정"
            ],
            "Admin GUI는 운영자용입니다."
        )

        # ==================== PART 9: TESTING ====================

        # Slide 61: Section Divider
        self.add_section_slide(9, "Testing Strategy", "테스트 전략")

        # Slide 62: Test Architecture
        self.add_content_slide(
            "🧪 테스트 아키텍처",
            [
                "## 테스트 피라미드",
                "- Integration Tests: End-to-End Order Flow",
                "- Component Tests: FMS, Navigation, AI",
                "- Unit Tests: Individual Functions & Classes",
                "",
                "## 컴포넌트별 테스트",
                "- FMS: OrderHandler, FleetController, StateMachine",
                "- Navigation: Waypoints, Path Planning, Recovery",
                "- AI: YOLO, Voice, Intent"
            ],
            "피라미드 형태의 테스트 아키텍처입니다."
        )

        # Slide 63: Skip Mode
        self.add_content_slide(
            "🧪 Skip Mode 아키텍처",
            [
                "## 문제",
                "- 물리적 로봇 없이 FMS 로직 테스트 필요",
                "",
                "## 해결: Skip Mode",
                "- skip_mode 파라미터로 활성화",
                "- 시뮬레이션: 설정된 시간 후 완료 처리",
                "- 실제 Navigation Action 대신 타이머 사용",
                "",
                "## 사용법",
                "- ros2 run fms fms_node --ros-args -p skip_mode:=true"
            ],
            "Skip Mode는 물리적 로봇 없이 FMS 로직을 테스트합니다."
        )

        # Slide 64: Test Coverage
        self.add_content_slide(
            "🧪 테스트 커버리지",
            [
                "## 전체 지표",
                "- Total Test Cases: 155",
                "- Code Coverage: 78%",
                "- Critical Path: 100%",
                "",
                "## 모듈별 커버리지",
                "- Order Handler: 45 tests, 92%",
                "- Fleet Controller: 38 tests, 85%",
                "- State Machine: 32 tests, 100%",
                "- TCP Server: 25 tests, 72%"
            ],
            "155개의 테스트 케이스로 78%의 커버리지를 달성했습니다."
        )

        # ==================== PART 10: PERFORMANCE ====================

        # Slide 65: Section Divider
        self.add_section_slide(10, "Performance & Scalability", "성능 및 확장성")

        # Slide 66: Performance Metrics
        self.add_two_column_slide(
            "📈 성능 지표",
            "Latency",
            [
                "주문 접수: <100ms",
                "Push 알림: <50ms",
                "Navigation Goal: <200ms",
                "YOLO 추론: <100ms",
                "Voice 처리: <2s"
            ],
            "Throughput",
            [
                "동시 주문: 10+",
                "로봇 수: 3 (확장 가능)",
                "초당 메시지: 100+",
                "평균 배달 시간: ~60s"
            ],
            "모든 지연 시간이 실시간 운영에 적합합니다."
        )

        # Slide 67: Database Schema
        self.add_content_slide(
            "🗄 데이터베이스 설계",
            [
                "## Core Tables",
                "- orders: 주문 정보, 상태, 타임스탬프",
                "- order_items: 주문 항목, 수량, 옵션",
                "- menu_items: 메뉴 카탈로그, 가격",
                "- robots: 로봇 상태, 위치, 배터리",
                "- delivery_history: 배달 기록, 분석용"
            ],
            "PostgreSQL 스키마입니다."
        )

        # Slide 68: Scalability
        self.add_two_column_slide(
            "📈 확장성 설계",
            "수평 확장",
            [
                "로봇 추가: Domain ID 할당만",
                "테이블 추가: Waypoint 설정만",
                "FMS 이중화: 상태 공유로 가능"
            ],
            "성능 개선",
            [
                "DB 캐싱: Redis 도입 예정",
                "비동기 처리: AsyncIO 적용",
                "로드 밸런싱: 로봇별 부하 분산"
            ],
            "수평 확장이 용이하도록 설계했습니다."
        )

        # ==================== PART 11: CONCLUSION ====================

        # Slide 69: Section Divider
        self.add_section_slide(11, "Conclusion", "결론 및 향후 계획")

        # Slide 70: Technical Lessons
        self.add_two_column_slide(
            "📚 기술적 교훈",
            "ROS2/DDS",
            [
                "WiFi에서 Multicast 불안정 → Unicast 필수",
                "Multi-Robot은 Multi-Domain으로 분리",
                "Domain Bridge로 통합 관리"
            ],
            "Architecture",
            [
                "Clean Architecture로 테스트 용이",
                "Dependency Injection으로 모듈화",
                "State Machine으로 Race Condition 방지"
            ],
            "프로젝트에서 얻은 기술적 교훈들입니다."
        )

        # Slide 71: Process Lessons
        self.add_content_slide(
            "📚 프로세스 교훈",
            [
                "## 조기 통합 테스트",
                "- 각 모듈 완성 전에 인터페이스부터 정의",
                "",
                "## 문서화의 중요성",
                "- API 문서, 아키텍처 다이어그램이 협업 효율 향상",
                "",
                "## 시뮬레이션 환경",
                "- Skip Mode 덕분에 HW 없이도 개발 진행",
                "",
                "## 로깅 체계",
                "- 구조화된 로그가 문제 해결 시간 단축"
            ],
            "개발 프로세스에서 배운 교훈들입니다."
        )

        # Slide 72: Future - Short Term
        self.add_content_slide(
            "🚀 향후 계획: 단기",
            [
                "## 3개월 내",
                "- 로봇 확장: 5대까지 지원",
                "- UI 개선: 웹 기반 Admin Dashboard",
                "- 음성 기능 강화: 다국어 지원",
                "- 에러 복구 자동화: 자동 재시도 로직",
                "",
                "## 6개월 내",
                "- Mobile App: 고객용 주문 앱",
                "- 분석 대시보드: 실시간 KPI 모니터링",
                "- 예측 모델: 주문량 예측"
            ],
            "단기 목표는 안정화와 기능 확장입니다."
        )

        # Slide 73: Future - Long Term
        self.add_two_column_slide(
            "🚀 향후 계획: 장기",
            "1년 내",
            [
                "다층 건물 네비게이션",
                "엘리베이터 연동",
                "자동 충전 시스템",
                "클라우드 FMS"
            ],
            "2년 내",
            [
                "멀티 매장 관리",
                "로봇 간 협업 (핸드오버)",
                "자율 재고 관리",
                "완전 무인 매장"
            ],
            "장기적으로는 완전 무인 매장을 목표로 합니다."
        )

        # Slide 74: Summary
        self.add_content_slide(
            "📌 요약",
            [
                "## Kitchmatics 달성 목표",
                "- ✅ 다중 로봇 관제: 3대 로봇 동시 운영",
                "- ✅ End-to-End 자동화: 주문~배달 전 과정",
                "- ✅ AI 품질 관리: YOLO 기반 음식 검사",
                "- ✅ 음성 주문: GPT-4 자연어 처리",
                "",
                "## 핵심 기술적 성과",
                "- WiFi 환경 DDS Discovery 문제 해결",
                "- Multi-Domain Architecture로 확장성 확보",
                "- Event-Driven State Machine으로 안정성 확보"
            ],
            "프로젝트의 핵심 성과를 요약했습니다."
        )

        # Slide 75: Thank You
        self.add_title_slide(
            "감사합니다",
            "Q&A"
        )

    def save(self, filename):
        self.prs.save(filename)
        print(f"PPTX saved: {filename}")


def main():
    gen = PPTXGenerator()
    gen.generate()
    gen.save("/home/gw/kitchmatics/roscamp-repo-1/presentation/kitchmatics_presentation_full.pptx")
    print(f"Total slides: {len(gen.prs.slides)}")


if __name__ == "__main__":
    main()
