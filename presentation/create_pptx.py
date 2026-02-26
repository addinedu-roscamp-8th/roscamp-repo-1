#!/usr/bin/env python3
"""
Kitchmatics 프로젝트 발표용 PPTX 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# 색상 정의
MAIN_COLOR = RGBColor(44, 62, 80)      # #2c3e50
ACCENT_COLOR = RGBColor(231, 76, 60)   # #e74c3c
SECONDARY_COLOR = RGBColor(52, 152, 219)  # #3498db
SUCCESS_COLOR = RGBColor(39, 174, 96)  # #27ae60
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 249, 250)

def set_slide_background(slide, color=WHITE):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """제목 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MAIN_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = SECONDARY_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, notes=""):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAIN_COLOR

    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = SECONDARY_COLOR
            p.space_before = Pt(15)
        elif item.startswith("-"):
            p.text = "  " + item
            p.font.size = Pt(18)
            p.font.color.rgb = MAIN_COLOR
            p.space_before = Pt(5)
        else:
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = MAIN_COLOR
            p.space_before = Pt(8)

    # 발표자 노트
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_diagram_slide(prs, title, diagram_text, notes=""):
    """다이어그램 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAIN_COLOR

    # 다이어그램 박스
    diagram_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.5))
    diagram_shape.fill.solid()
    diagram_shape.fill.fore_color.rgb = LIGHT_GRAY
    diagram_shape.line.color.rgb = RGBColor(200, 200, 200)

    # 다이어그램 텍스트
    tf = diagram_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = MAIN_COLOR

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, notes=""):
    """2열 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAIN_COLOR

    # 왼쪽 컬럼 제목
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # 왼쪽 컬럼 내용
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item if not item.startswith("-") else "  " + item
        p.font.size = Pt(16)
        p.font.color.rgb = MAIN_COLOR
        p.space_before = Pt(5)

    # 오른쪽 컬럼 제목
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # 오른쪽 컬럼 내용
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item if not item.startswith("-") else "  " + item
        p.font.size = Pt(16)
        p.font.color.rgb = MAIN_COLOR
        p.space_before = Pt(5)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_problem_solution_slide(prs, title, problem_items, solution_items, notes=""):
    """문제-해결 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAIN_COLOR

    # 문제 박스
    problem_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(4.5), Inches(2.8))
    problem_shape.fill.solid()
    problem_shape.fill.fore_color.rgb = RGBColor(255, 245, 245)
    problem_shape.line.color.rgb = ACCENT_COLOR
    problem_shape.line.width = Pt(2)

    tf = problem_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚨 문제 상황"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    for item in problem_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = MAIN_COLOR
        p.space_before = Pt(5)

    # 해결 박스
    solution_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.5), Inches(2.8))
    solution_shape.fill.solid()
    solution_shape.fill.fore_color.rgb = RGBColor(240, 255, 244)
    solution_shape.line.color.rgb = SUCCESS_COLOR
    solution_shape.line.width = Pt(2)

    tf = solution_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ 해결책"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR

    for item in solution_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = MAIN_COLOR
        p.space_before = Pt(5)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===========================================
    # Slide 1: 표지
    # ===========================================
    slide = add_title_slide(prs,
        "🍽️ Kitchmatics",
        "자율주행 서빙 로봇 기반 자동화 식당 운영 시스템")

    # 키워드 추가
    keyword_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    tf = keyword_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ROS2  |  Multi-Robot  |  AI Vision  |  Voice Order"
    p.font.size = Pt(18)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(10), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026.02"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """안녕하세요. Kitchmatics 프로젝트를 소개하겠습니다.
Kitchmatics는 자율주행 서빙 로봇을 활용한 완전 자동화 식당 운영 시스템입니다.
ROS2 기반 멀티 로봇 제어, AI 비전을 통한 음식 검수, 음성 주문 시스템까지 풀스택 자동화를 구현했습니다."""

    # ===========================================
    # Slide 2: 목차
    # ===========================================
    add_two_column_slide(prs, "📋 목차",
        "Part 1",
        ["1. 프로젝트 개요",
         "2. 시스템 아키텍처",
         "3. 핵심 기술 - FMS",
         "4. 핵심 기술 - Robot Arm",
         "5. 핵심 기술 - AI Server",
         "6. 핵심 기술 - GUI"],
        "Part 2",
        ["7. 통신 아키텍처",
         "8. 문제 해결 사례",
         "9. 테스트 전략",
         "10. 데모 시나리오",
         "11. 성과 및 결론",
         "12. Q&A"],
        notes="""발표는 총 12개 섹션으로 구성됩니다.
프로젝트 개요부터 시작해서 시스템 아키텍처, 핵심 기술들(FMS, 로봇팔, AI 서버, GUI),
통신 아키텍처, 그리고 가장 중요한 문제 해결 사례를 말씀드리겠습니다.
마지막으로 테스트 전략과 성과를 공유하고 Q&A 시간을 갖겠습니다.""")

    # ===========================================
    # Slide 3: 프로젝트 개요
    # ===========================================
    add_two_column_slide(prs, "1. 프로젝트 개요",
        "🎯 목표",
        ["완전 자동화 식당 운영",
         "다중 로봇(3대) 협업 서빙",
         "음성 주문 인터페이스",
         "AI 비전 품질 검수",
         "",
         "📊 프로젝트 규모",
         "- 10,000+ Lines of Code",
         "- 155 Test Cases",
         "- 14+ 문서 파일"],
        "🔧 핵심 구성요소",
        ["FMS (Fleet Management System)",
         "- 로봇 함대 관리 중추",
         "",
         "Mobile Robot x3 (pinky1/2/3)",
         "- Turtlebot3 기반 자율주행",
         "",
         "Robot Arm x2",
         "- MyCobot 280 샌드위치 조립",
         "",
         "AI Server",
         "- YOLO 검수 + 음성 주문"],
        notes="""프로젝트의 목표는 고객 주문부터 음식 조리, 서빙, 수령 확인까지 전 과정을 자동화하는 것입니다.

특히 3대의 서빙 로봇이 동시에 협업하는 다중 로봇 시스템,
음성으로 주문할 수 있는 인터페이스,
그리고 AI 비전으로 음식 품질을 검수하는 기능을 구현했습니다.

전체 코드는 1만 라인 이상이며, 155개의 테스트 케이스로 검증되었습니다.""")

    # ===========================================
    # Slide 4: 시스템 아키텍처
    # ===========================================
    add_diagram_slide(prs, "2. 시스템 아키텍처",
        """
    ┌─────────────────────────────────────────────────────────┐
    │                    Customer Zone                         │
    │   [Kiosk T1]      [Kiosk T2]      [Kiosk T3]           │
    └───────────────────────┬─────────────────────────────────┘
                            │ TCP/JSON (Port 9000)
    ┌───────────────────────▼─────────────────────────────────┐
    │              Master PC (192.168.1.3)                     │
    │   ┌─────────────────────────────────────────────────┐   │
    │   │              FMS Node (Domain 25)                │   │
    │   │   [OrderHandler] [FleetCtrl] [PathPlanner]      │   │
    │   └─────────────────────┬───────────────────────────┘   │
    │                         │ Domain Bridge                  │
    │           ┌─────────────┼─────────────┐                 │
    │           ▼             ▼             ▼                 │
    │   [pinky1 D:11]  [pinky2 D:12]  [pinky3 D:13]          │
    └─────────────────────────────────────────────────────────┘

    기술 스택: ROS2 Humble | CycloneDDS | Multi-Domain | PostgreSQL
        """,
        notes="""전체 시스템 아키텍처입니다.

상단은 고객 영역으로, 각 테이블에 키오스크가 있습니다.
키오스크는 TCP/JSON으로 FMS와 통신합니다.

중앙의 Master PC에서 FMS가 실행되며,
OrderHandler가 주문을 처리하고,
FleetController가 로봇 배차를 담당하고,
PathPlanner가 경로를 계획합니다.

핵심은 Multi-Domain 아키텍처입니다.
FMS는 Domain 25에서 실행되고,
각 로봇은 Domain 11, 12, 13에서 독립적으로 실행됩니다.
Domain Bridge가 이들을 연결합니다.""")

    # ===========================================
    # Slide 5: 기술 스택
    # ===========================================
    add_content_slide(prs, "기술 스택",
        ["## Robot Control",
         "- ROS2 Humble, Nav2, AMCL, CycloneDDS",
         "",
         "## AI / ML",
         "- YOLOv8 (음식 검수), OpenAI Whisper (STT)",
         "- GPT-4 Function Calling (주문 파싱)",
         "",
         "## Backend",
         "- Python 3.12, FastAPI, SQLAlchemy, PostgreSQL",
         "",
         "## Frontend",
         "- PyQt5, TCP Socket, Qt Designer",
         "",
         "## Hardware",
         "- Turtlebot3 (서빙봇), MyCobot 280 (로봇팔), Jetson Nano"],
        notes="""사용한 기술 스택입니다.

로봇 제어에는 ROS2 Humble, Nav2 네비게이션 스택, AMCL 로컬라이제이션을 사용했고,
DDS는 CycloneDDS를 선택했습니다.

AI에는 YOLOv8로 음식 검수, OpenAI Whisper로 음성 인식,
GPT-4 Function Calling으로 자연어 주문 파싱을 구현했습니다.""")

    # ===========================================
    # Slide 6: FMS 아키텍처
    # ===========================================
    add_two_column_slide(prs, "3. 핵심 기술 - FMS (Fleet Management System)",
        "📦 주요 모듈 (10,225 LoC)",
        ["fms_node.py (2,361 lines)",
         "- 메인 오케스트레이터",
         "",
         "collision_avoidance (1,297 lines)",
         "- 다중 로봇 충돌 회피",
         "",
         "task_scheduler (693 lines)",
         "- 작업 스케줄링, 배차",
         "",
         "order_handler (750 lines)",
         "- 주문 처리, 상태 머신"],
        "🎯 설계 원칙",
        ["Clean Architecture",
         "- Domain → Application → Infrastructure",
         "",
         "SOLID Principles",
         "- Dependency Inversion",
         "- Single Responsibility",
         "",
         "Callback-based Design",
         "- 레이어 간 느슨한 결합",
         "",
         "Event-Driven",
         "- 비동기 상태 전환"],
        notes="""FMS는 전체 시스템의 두뇌 역할을 합니다.

총 10,225 라인의 코드로 구성되어 있으며,
가장 큰 fms_node.py가 2,361 라인입니다.

충돌 회피 모듈이 1,297 라인으로 두 번째로 큰데,
이는 다중 로봇 환경에서 충돌 방지가 얼마나 복잡한지를 보여줍니다.""")

    # ===========================================
    # Slide 7: 주문 상태 머신
    # ===========================================
    add_diagram_slide(prs, "주문 상태 머신 (Order State Machine)",
        """
    ┌─────────────┐
    │  RECEIVED   │  ← 주문 접수 (GUI → FMS)
    └──────┬──────┘
           │ 조리 명령 + 로봇 출발
           ▼
    ┌─────────────┐
    │   COOKING   │  ← 로봇팔 조리 중 + pinky가 point13으로 이동
    └──────┬──────┘
           │ 로봇 도착 (point13)
           ▼
    ┌─────────────┐
    │   LOADING   │  ← 조리 완료 대기 (이벤트 기반 동기화)
    └──────┬──────┘
           │ LoadingComplete 수신
           ▼
    ┌─────────────┐
    │   LOADED    │  ← 음식 적재 완료
    └──────┬──────┘
           │ navigate_robot('tableN')
           ▼
    ┌─────────────┐
    │ DELIVERING  │  ← Nav2로 테이블 이동
    └──────┬──────┘
           │ 테이블 도착 → Push 알림
           ▼
    ┌─────────────┐
    │  COMPLETED  │  ← 수령 확인 → 주차장 복귀
    └─────────────┘
        """,
        notes="""주문 상태 머신입니다.

주문이 접수되면 RECEIVED 상태에서 시작합니다.
동시에 로봇팔에 조리 명령을 보내고, 서빙 로봇은 픽업 포인트로 이동합니다.

여기서 중요한 문제가 발생했습니다.
로봇이 픽업 포인트에 먼저 도착했는데 조리가 안 끝났을 때,
단순히 타이머로 기다리면 조리 시간이 가변적이라 문제가 됩니다.

해결책은 이벤트 기반 전환입니다.
조리 완료 토픽을 구독하고, 두 조건이 모두 충족될 때만 다음 단계로 넘어갑니다.""")

    # ===========================================
    # Slide 8: Robot Arm
    # ===========================================
    add_two_column_slide(prs, "4. 핵심 기술 - Robot Arm",
        "🦾 시스템 구성",
        ["Sandwich Arm (ARM_1)",
         "- 재료 집기, 빵 위에 적재",
         "",
         "Sauce Arm (ARM_2)",
         "- 소스 도포",
         "",
         "Clean Architecture 적용",
         "- Domain/Application/Infra 분리",
         "",
         "📡 ROS2 토픽",
         "/arm/command (FMS → Arm)",
         "/arm/status (Arm → FMS)"],
        "📊 성능 및 인터페이스",
        ["명령 수신 → 실행: <10ms",
         "상태 발행 주기: 0.5초",
         "처리량: 100+ msg/s",
         "",
         "JSON 메시지 형식:",
         "{",
         "  \"job_id\": \"JOB-001\",",
         "  \"operation\": \"START\",",
         "  \"order\": {...}",
         "}"],
        notes="""로봇팔 시스템입니다.

2대의 MyCobot 280을 사용하는데,
Sandwich Arm은 재료를 집어서 빵 위에 올리고,
Sauce Arm은 소스를 도포합니다.

FMS와의 통신은 JSON 기반 ROS2 토픽을 사용합니다.
명령 처리 레이턴시는 10ms 이내입니다.""")

    # ===========================================
    # Slide 9: AI Server - YOLO
    # ===========================================
    add_two_column_slide(prs, "5. 핵심 기술 - AI Server (YOLO)",
        "🎯 목적 및 구성",
        ["샌드위치 품질 검수",
         "- 재료 누락 감지",
         "- 조립 상태 확인",
         "",
         "🔧 기술 스택",
         "모델: YOLOv8 Custom (best.pt)",
         "클래스: m1, m2, m3 (메뉴별)",
         "서버: Flask (Port 5001)",
         "소스: JetBot 카메라 스트림"],
        "📡 API Endpoints",
        ["GET /analyze/image",
         "- 스냅샷 1장 분석",
         "",
         "GET /analyze/video",
         "- 영상 N프레임 분석",
         "",
         "GET /view",
         "- 실시간 웹 뷰",
         "",
         "응답: JSON (detections, confidence, bbox)"],
        notes="""YOLO 기반 AI 분석 서버입니다.

목적은 샌드위치 품질 검수입니다.
로봇팔이 조립한 샌드위치를 카메라로 촬영하고,
YOLO 모델이 재료 누락이나 조립 불량을 감지합니다.""")

    # ===========================================
    # Slide 10: AI Server - Voice
    # ===========================================
    add_content_slide(prs, "핵심 기술 - AI Server (Voice)",
        ["## 🎤 음성 주문 파이프라인",
         "음성 입력 → STT (Whisper) → Intent 분석 → Function Call → TTS 응답",
         "",
         "## 🔧 핵심 기술",
         "- STT: OpenAI Whisper API",
         "- Intent: GPT-4 Function Calling",
         "- TTS: OpenAI TTS API",
         "- Wakeword: 주문 시작 의도 감지",
         "",
         "## 📡 API Endpoints",
         "- POST /stt/ - 음성 → 텍스트",
         "- POST /tts/ - 텍스트 → 음성",
         "- POST /wakeword/check - 주문 의도 감지",
         "- POST /pipeline/run - 전체 파이프라인",
         "",
         "## 💬 예시",
         "Input: \"햄치즈샌드위치 2개 주세요\"",
         "Output: {menu_id: \"M001\", quantity: 2}"],
        notes="""음성 주문 시스템입니다.

파이프라인은 다음과 같습니다.
고객이 음성으로 주문하면, Whisper가 텍스트로 변환합니다.
GPT-4 Function Calling이 자연어를 구조화된 주문으로 파싱합니다.""")

    # ===========================================
    # Slide 11: GUI System
    # ===========================================
    add_two_column_slide(prs, "6. 핵심 기술 - GUI System",
        "🖥️ Customer GUI (키오스크)",
        ["주문 시작 화면",
         "- 터치로 주문 시작",
         "",
         "메뉴 선택 화면",
         "- 메뉴 리스트, 장바구니, 수량",
         "",
         "주문 확인 화면",
         "- 영수증 형태 표시",
         "",
         "수령 확인 화면",
         "- 도착 알림, 수령 버튼",
         "",
         "음성 피드백 위젯",
         "- 파형 애니메이션"],
        "👨‍💼 Admin GUI (관리자)",
        ["주문 대시보드",
         "- 실시간 주문 모니터링",
         "",
         "조리 모니터",
         "- 조리 상태, 검수 결과",
         "",
         "레시피/재고 관리",
         "- CRUD, 알림",
         "",
         "🚗 Fleet 모니터 (NEW!)",
         "- 로봇 3대 실시간 상태",
         "- 배터리, 작업 표시"],
        notes="""GUI는 Customer용과 Admin용 두 가지입니다.

Customer GUI는 테이블에 설치되는 키오스크입니다.
Admin GUI의 Fleet 모니터 화면에서 3대 로봇의 위치와 상태를 한눈에 볼 수 있습니다.""")

    # ===========================================
    # Slide 12: 통신 아키텍처
    # ===========================================
    add_diagram_slide(prs, "7. 통신 아키텍처",
        """
    ┌─────────────────────────────────────────────────────────┐
    │               Network: kitchmatics WiFi (192.168.1.x)    │
    │                                                          │
    │    [Customer GUI] ──TCP:9000──┐                         │
    │    [Admin GUI] ────TCP:9999───┤                         │
    │                               ▼                          │
    │                    ┌─────────────────┐                   │
    │                    │   FMS Node      │                   │
    │                    │  192.168.1.3    │                   │
    │                    │   Domain 25     │                   │
    │                    └────────┬────────┘                   │
    │                             │ Domain Bridge              │
    │           ┌─────────────────┼─────────────────┐         │
    │           ▼                 ▼                 ▼         │
    │    [pinky1]          [pinky2]          [pinky3]         │
    │   192.168.1.7       192.168.1.6      192.168.1.11       │
    │    Domain 11         Domain 12        Domain 13         │
    │                                                          │
    │    [Robot Arm]       [AI Server]                        │
    │   192.168.1.4       192.168.0.27                        │
    │    Domain 20         REST API                           │
    └─────────────────────────────────────────────────────────┘

    통신 방식: TCP/JSON (GUI↔FMS) | ROS2/DDS (FMS↔Robot) | REST (FMS↔AI)
        """,
        notes="""통신 아키텍처입니다.

세 가지 통신 방식을 사용합니다.
첫째, GUI와 FMS는 TCP/JSON으로 통신합니다.
둘째, FMS와 로봇은 ROS2/DDS로 통신합니다.
셋째, AI Server와는 REST API로 통신합니다.""")

    # ===========================================
    # Slide 13: 문제 해결 1 - DDS Discovery
    # ===========================================
    add_problem_solution_slide(prs, "8. 문제 해결 ① DDS Discovery 실패",
        ["SSH는 되는데 ROS2 토픽이 안 보임",
         "ros2 topic list 결과가 비어있음",
         "로봇은 정상 동작 중",
         "",
         "🔍 원인:",
         "ROS2 기본: Multicast UDP 사용",
         "WiFi 공유기가 Multicast 차단",
         "노드 간 Discovery 불가능"],
        ["Unicast Discovery 설정",
         "",
         "cyclonedds_main.xml:",
         "<Peers>",
         "  <Peer address=\"192.168.1.7\"/>",
         "  <Peer address=\"192.168.1.6\"/>",
         "</Peers>",
         "",
         "📈 결과:",
         "Peer IP 직접 지정으로 안정적 통신",
         "WiFi 환경에서도 100% 연결 성공"],
        notes="""첫 번째 문제 해결 사례입니다.

문제는 SSH로 로봇에 접속은 되는데, ROS2 토픽이 전혀 보이지 않는 것이었습니다.

원인을 분석해보니, ROS2가 기본으로 Multicast UDP를 사용하는데,
WiFi 공유기가 Multicast 패킷을 차단하고 있었습니다.

해결책은 CycloneDDS의 Unicast Discovery 기능을 사용하는 것입니다.""")

    # ===========================================
    # Slide 14: 문제 해결 2 - Multi-Domain
    # ===========================================
    add_problem_solution_slide(prs, "문제 해결 ② Multi-Domain 토픽 충돌",
        ["로봇 3대가 같은 토픽명 사용",
         "/cmd_vel 명령이 모든 로봇에 전달",
         "개별 제어 불가능",
         "",
         "🔍 원인:",
         "동일 Domain에서 동일 토픽명 충돌",
         "Namespace만으로는 불충분",
         "Nav2 Action도 충돌"],
        ["Multi-Domain + Bridge",
         "",
         "Domain 할당:",
         "FMS: Domain 25",
         "pinky1: Domain 11",
         "pinky2: Domain 12",
         "pinky3: Domain 13",
         "",
         "📈 결과:",
         "로봇별 완전 격리",
         "필요한 토픽만 선택적 Bridge"],
        notes="""두 번째 문제 해결 사례입니다.

문제는 로봇 3대가 모두 같은 토픽명을 사용해서 충돌이 발생한 것입니다.

해결책은 Multi-Domain 아키텍처입니다.
각 로봇에 별도의 Domain ID를 부여하여 완전히 격리합니다.""")

    # ===========================================
    # Slide 15: 문제 해결 3 - 상태 동기화
    # ===========================================
    add_problem_solution_slide(prs, "문제 해결 ③ 로봇-조리 상태 동기화",
        ["로봇이 픽업 포인트 도착",
         "조리가 아직 안 끝남",
         "언제 출발해야 하나?",
         "",
         "❌ 실패한 시도:",
         "고정 타이머 (30초 대기)",
         "→ 조리 시간이 가변적이라 실패",
         "폴링 방식",
         "→ 지연 발생, 리소스 낭비"],
        ["Event-Driven + 조건부 전환",
         "",
         "check_can_proceed():",
         "  robot_at_pickup = (state == AT_POINT13)",
         "  cooking_done = loading_complete_received",
         "",
         "  if robot_at_pickup AND cooking_done:",
         "      transition_to(DELIVERING)",
         "",
         "📈 결과:",
         "이벤트 기반 상태 전환",
         "두 조건 AND 충족 시만 진행"],
        notes="""세 번째 문제입니다.

로봇이 픽업 포인트에 도착했는데 음식이 아직 준비 안 됐을 때,
언제 출발해야 하는지의 문제입니다.

최종 해결책은 이벤트 기반 설계입니다.
두 조건이 AND로 충족될 때만 다음 단계로 전환합니다.""")

    # ===========================================
    # Slide 16: 문제 해결 4 - 에러 복구
    # ===========================================
    add_two_column_slide(prs, "문제 해결 ④ 서빙 오류 대응 시스템",
        "🔍 오류 유형",
        ["NAV_FAILED - 경로 실패 (즉시)",
         "COMM_LOST - 통신 끊김 (5초)",
         "LOW_BATTERY - 배터리 부족 (20V)",
         "TIMEOUT - 작업 타임아웃 (60/120초)",
         "OBSTACLE - 장애물 (3회 실패)",
         "",
         "📡 알림 흐름",
         "ErrorDetector (0.5Hz 모니터링)",
         "→ /fms/error_alert 발행",
         "→ Main Server (TCP broadcast)",
         "→ Admin GUI (팝업 알림)"],
        "⚙️ 복구 옵션",
        ["RETRY - 현재 위치에서 재시도",
         "",
         "RETURN_HOME - 주차장으로 강제 복귀",
         "",
         "EMERGENCY_STOP - 긴급 정지",
         "",
         "CLEAR_ERROR - 에러 상태만 해제",
         "",
         "운영자가 상황에 맞게 선택",
         "Admin GUI에서 원클릭 실행"],
        notes="""네 번째는 오류 대응 시스템입니다.

서빙 중 발생할 수 있는 오류 유형을 정의했습니다.
ErrorDetector가 0.5Hz로 모니터링하다가 오류를 감지하면 알림을 발송합니다.""")

    # ===========================================
    # Slide 17: 테스트 전략
    # ===========================================
    add_two_column_slide(prs, "9. 테스트 전략",
        "📊 테스트 커버리지",
        ["Unit Tests (~80개)",
         "- test_fms_unit.py",
         "- Task, TaskManager, RobotState",
         "",
         "Integration Tests (~35개)",
         "- test_multi_robot.py",
         "- 다중 로봇 시나리오",
         "",
         "E2E Tests (~40개)",
         "- test_e2e_skip_mode.py",
         "- 전체 플로우 검증",
         "",
         "Total: 155 Tests"],
        "⚡ Skip Mode",
        ["목적:",
         "- 실제 로봇/로봇팔 없이 테스트",
         "- 외부 의존성 제거",
         "- 빠른 개발 사이클",
         "",
         "동작 방식:",
         "- Precision parking: 2초 Mock",
         "- Food loading: 3초 Mock",
         "- 상태 전환은 실제와 동일",
         "",
         "ros2 run fms fms_node \\",
         "  --ros-args -p skip_mode:=true"],
        notes="""테스트 전략입니다.

총 155개 테스트를 작성했습니다.
핵심은 Skip Mode입니다. 실제 로봇이나 로봇팔 없이도 전체 시스템을 테스트할 수 있습니다.""")

    # ===========================================
    # Slide 18: 데모 시나리오
    # ===========================================
    add_content_slide(prs, "10. 데모 시나리오",
        ["## 1️⃣ 주문 접수",
         "고객이 키오스크에서 \"햄치즈샌드위치 2개\" 음성 주문",
         "",
         "## 2️⃣ 조리 시작",
         "FMS가 Robot Arm에 조리 명령, pinky1 픽업 포인트로 출발",
         "",
         "## 3️⃣ 품질 검수",
         "AI Server가 YOLO로 샌드위치 검수, 정상 판정",
         "",
         "## 4️⃣ 음식 적재",
         "pinky1 도착, 로봇팔이 샌드위치를 로봇에 적재",
         "",
         "## 5️⃣ 테이블 서빙",
         "pinky1이 테이블로 이동 (Nav2 네비게이션)",
         "",
         "## 6️⃣ 수령 완료",
         "키오스크 푸시 알림 → 고객 수령 버튼 클릭 → 주차장 복귀",
         "",
         "⏱️ 전체 소요시간: 약 3-5분"],
        notes="""데모 시나리오입니다.

고객이 키오스크에서 음성으로 주문합니다.
FMS가 주문을 받아서 로봇팔과 서빙 로봇을 동시에 제어합니다.
전체 소요시간은 약 3-5분입니다.""")

    # ===========================================
    # Slide 19: 성능 지표
    # ===========================================
    add_two_column_slide(prs, "성능 지표",
        "⚡ 레이턴시",
        ["주문 접수: <100ms",
         "푸시 알림: <50ms",
         "주문→서빙 완료: 3-5분",
         "",
         "🤖 로봇 성능",
         "최대 속도: 0.26 m/s",
         "배터리: ~2시간 연속",
         "위치 정확도: ±5cm (AMCL)",
         "Heartbeat: 0.5초"],
        "💻 시스템 성능",
        ["FMS CPU 사용: ~5% (idle)",
         "FMS 메모리: ~200MB",
         "DB 연결 풀: 10 connections",
         "동시 주문 처리: 3+",
         "",
         "📊 코드 품질",
         "10,225 Lines of Code",
         "155 Test Cases",
         "14+ 문서 파일"],
        notes="""성능 지표입니다.

주문 접수 레이턴시는 100ms 이내입니다.
푸시 알림 레이턴시는 50ms 이내입니다.
로봇은 AMCL 기준 ±5cm 정확도로 위치를 추정합니다.""")

    # ===========================================
    # Slide 20: 배운 점
    # ===========================================
    add_two_column_slide(prs, "11. 배운 점 / 교훈",
        "✅ 잘한 점",
        ["Multi-Domain 아키텍처",
         "- 로봇 격리로 충돌 문제 해결",
         "",
         "Skip Mode 도입",
         "- 개발 속도 3배 향상",
         "",
         "Event-Driven 설계",
         "- 상태 동기화 문제 해결",
         "",
         "Unicast DDS 설정",
         "- WiFi 환경 안정화"],
        "📚 개선 포인트",
        ["초기 아키텍처 설계",
         "- Multi-Domain을 처음부터 고려했으면...",
         "",
         "네트워크 환경 테스트",
         "- 실제 WiFi 환경에서 조기 테스트 필요",
         "",
         "에러 복구 자동화",
         "- 현재 수동 개입 필요",
         "",
         "부하 테스트",
         "- 100+ 동시 주문 시나리오 미검증"],
        notes="""배운 점과 교훈입니다.

잘한 점은 Multi-Domain 아키텍처, Skip Mode, Event-Driven 설계입니다.
개선 포인트로는 초기 설계 시 Multi-Domain을 고려했으면 좋았을 것입니다.""")

    # ===========================================
    # Slide 21: 향후 계획
    # ===========================================
    add_two_column_slide(prs, "향후 계획",
        "🔜 단기 계획",
        ["자동 에러 복구",
         "- 특정 오류 자동 재시도",
         "",
         "배터리 기반 스케줄링",
         "- 배터리 잔량 고려한 배차",
         "",
         "충돌 회피 고도화",
         "- 동적 장애물 대응",
         "",
         "통계 대시보드",
         "- 주문/로봇 분석"],
        "🚀 장기 계획",
        ["로봇 확장",
         "- 5대 이상 동시 운영",
         "",
         "멀티 플로어",
         "- 엘리베이터 연동",
         "",
         "메뉴 다양화",
         "- 다양한 레시피 지원",
         "",
         "프랜차이즈 확장",
         "- 클라우드 기반 관제"],
        notes="""향후 계획입니다.

단기적으로는 자동 에러 복구, 배터리 기반 스케줄링을 구현할 예정입니다.
장기적으로는 로봇 확장, 멀티 플로어, 프랜차이즈 확장을 고려하고 있습니다.""")

    # ===========================================
    # Slide 22: 결론
    # ===========================================
    add_content_slide(prs, "결론",
        ["## 🎯 Kitchmatics 핵심 성과",
         "",
         "## 🤖 Multi-Robot System",
         "3대 동시 협업, Domain 격리로 충돌 없는 운영",
         "",
         "## 🧠 AI Integration",
         "YOLO 품질 검수, 음성 주문으로 사용자 경험 향상",
         "",
         "## ⚡ Reliability",
         "155 테스트, Skip Mode로 높은 신뢰성 확보",
         "",
         "",
         "\"주문부터 서빙까지 완전 자동화된 식당 운영 시스템\"",
         "",
         "10,225 LoC  |  155 Tests  |  3 Robots  |  Production Ready"],
        notes="""결론입니다.

Kitchmatics의 핵심 성과는 세 가지입니다.
첫째, 3대의 로봇이 동시에 협업하는 Multi-Robot 시스템
둘째, AI 통합 (YOLO 검수, 음성 주문)
셋째, 높은 신뢰성 (155개 테스트, Skip Mode)

감사합니다.""")

    # ===========================================
    # Slide 23: Q&A
    # ===========================================
    slide = add_title_slide(prs, "12. Q&A", "감사합니다")

    # 키워드 추가
    keyword_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf = keyword_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ROS2  |  Multi-Robot  |  AI Vision  |  Voice Order  |  Clean Architecture"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """질문 있으시면 받겠습니다.

발표를 들어주셔서 감사합니다.
궁금하신 점이 있으시면 편하게 질문해 주세요."""

    # ===========================================
    # 저장
    # ===========================================
    output_path = "/home/gw/kitchmatics/roscamp-repo-1/presentation/kitchmatics_presentation.pptx"
    prs.save(output_path)
    print(f"✅ PPTX 파일 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
